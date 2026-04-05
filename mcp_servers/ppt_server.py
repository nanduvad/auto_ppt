from pptx import Presentation

class PPTServer:
    def __init__(self):
        self.prs = Presentation()

    def create_presentation(self):
        self.prs = Presentation()
        return "Presentation created"

    def add_slide(self, title, bullets):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        content = slide.placeholders[1]

        title_placeholder.text = title

        tf = content.text_frame
        tf.clear()

        for point in bullets:
            p = tf.add_paragraph()
            p.text = point

        return f"Slide added: {title}"

    def save_presentation(self, filename="output.pptx"):
        self.prs.save(filename)
        return f"Saved as {filename}"