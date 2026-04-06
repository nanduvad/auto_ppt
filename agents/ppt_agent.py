from agents.base_agent import BaseAgent
from pptx import Presentation
from pptx.util import Pt, Emu, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches, Emu
import json


# ---------------------------------
# THEME DEFINITIONS
# ---------------------------------
THEMES = {
    "modern": {
        # ── background ──────────────────────────────────────────────────────
        "bg":             RGBColor(15,  23,  42),   # slate-900 deep navy
        # ── title slide ─────────────────────────────────────────────────────
        "slide_title":    RGBColor(255, 220, 100),  # warm gold  ← title slide topic color
        "title_bg":       RGBColor(20,  30,  55),   # darker navy hero panel
        "subtitle":       RGBColor(148, 163, 184),  # slate-400 tagline
        # ── content slide heading ────────────────────────────────────────────
        "title":          RGBColor(99,  210, 255),  # sky-blue  ← slide heading color
        "heading_bg":     RGBColor(22,  36,  71),   # navy band behind heading
        # ── paragraph / body text ───────────────────────────────────────────
        "text":           RGBColor(226, 232, 240),  # slate-200  ← paragraph color
        # ── accents & cards ─────────────────────────────────────────────────
        "accent":         RGBColor(56,  189, 248),  # sky-400
        "card_bg":        RGBColor(30,  41,  59),   # slate-800
        "card_border":    RGBColor(51,  65,  85),   # slate-700
        # ── typography ──────────────────────────────────────────────────────
        "title_font":     "Calibri",
        "body_font":      "Calibri",
        "title_size":     Pt(34),
        "body_size":      Pt(15),
    },
    "dark": {
        # ── background ──────────────────────────────────────────────────────
        "bg":             RGBColor(10,  10,  10),   # true black
        # ── title slide ─────────────────────────────────────────────────────
        "slide_title":    RGBColor(255, 255, 255),  # pure white  ← title slide topic color
        "title_bg":       RGBColor(15,  10,  30),   # near-black violet tint
        "subtitle":       RGBColor(161, 161, 170),  # zinc-400 tagline
        # ── content slide heading ────────────────────────────────────────────
        "title":          RGBColor(196, 167, 255),  # soft violet  ← slide heading color
        "heading_bg":     RGBColor(20,  14,  40),   # dark violet band
        # ── paragraph ───────────────────────────────────────────────────────
        "text":           RGBColor(228, 228, 231),  # zinc-200  ← paragraph color
        # ── accents & cards ─────────────────────────────────────────────────
        "accent":         RGBColor(167, 139, 250),  # violet-400
        "card_bg":        RGBColor(24,  24,  27),   # zinc-900
        "card_border":    RGBColor(63,  63,  70),   # zinc-700
        # ── typography ──────────────────────────────────────────────────────
        "title_font":     "Arial Black",
        "body_font":      "Calibri",
        "title_size":     Pt(34),
        "body_size":      Pt(15),
    },
    "minimal": {
        # ── background ──────────────────────────────────────────────────────
        "bg":             RGBColor(250, 250, 249),  # warm off-white
        # ── title slide ─────────────────────────────────────────────────────
        "slide_title":    RGBColor(255, 255, 255),  # white on dark hero  ← title slide topic color
        "title_bg":       RGBColor(28,  25,  23),   # stone-900 dark hero panel
        "subtitle":       RGBColor(180, 170, 160),  # muted warm grey tagline
        # ── content slide heading ────────────────────────────────────────────
        "title":          RGBColor(194, 65,  12),   # burnt orange  ← slide heading color
        "heading_bg":     RGBColor(255, 245, 235),  # light peach band
        # ── paragraph ───────────────────────────────────────────────────────
        "text":           RGBColor(41,  37,  36),   # stone-800  ← paragraph color
        # ── accents & cards ─────────────────────────────────────────────────
        "accent":         RGBColor(234, 88,  12),   # orange-600
        "card_bg":        RGBColor(243, 240, 235),  # warm grey card
        "card_border":    RGBColor(214, 211, 209),  # stone-300
        # ── typography ──────────────────────────────────────────────────────
        "title_font":     "Georgia",
        "body_font":      "Calibri",
        "title_size":     Pt(36),
        "body_size":      Pt(15),
    },
    "corporate": {
        # ── background ──────────────────────────────────────────────────────
        "bg":             RGBColor(235, 243, 250),  # pale ice blue
        # ── title slide ─────────────────────────────────────────────────────
        "slide_title":    RGBColor(255, 220, 80),   # bright gold  ← title slide topic color
        "title_bg":       RGBColor(2,   62,  84),   # deep teal hero panel
        "subtitle":       RGBColor(160, 200, 220),  # light teal tagline
        # ── content slide heading ────────────────────────────────────────────
        "title":          RGBColor(245, 168, 35),   # golden amber  ← slide heading color
        "heading_bg":     RGBColor(2,   62,  84),   # deep teal band
        # ── paragraph ───────────────────────────────────────────────────────
        "text":           RGBColor(15,  40,  55),   # dark teal-ink  ← paragraph color
        # ── accents & cards ─────────────────────────────────────────────────
        "accent":         RGBColor(245, 168, 35),   # amber
        "card_bg":        RGBColor(214, 232, 245),  # light blue card
        "card_border":    RGBColor(180, 210, 230),
        # ── typography ──────────────────────────────────────────────────────
        "title_font":     "Cambria",
        "body_font":      "Calibri",
        "title_size":     Pt(34),
        "body_size":      Pt(15),
    },
}

# Slide canvas dimensions (standard 16:9 widescreen)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


class PPTAgent(BaseAgent):

    def __init__(self):
        super().__init__(
            name="ppt_generator",
            role="Creates professional PowerPoint slides"
        )

    # ---------------------------------
    # APPLY THEME TO SLIDE
    # ---------------------------------
    def apply_theme(self, slide, theme_style):
        """Fill slide background with theme base colour."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme_style["bg"]

    # ---------------------------------
    # STYLE TITLE
    # ---------------------------------
    def style_title(self, text_frame, theme_style):
        """Style the primary title text frame."""
        para = text_frame.paragraphs[0]
        font = para.font
        font.name  = theme_style["title_font"]
        font.size  = Pt(32)
        font.bold  = True
        font.color.rgb = theme_style["title"]

    # ---------------------------------
    # STYLE BULLETS
    # ---------------------------------
    def style_bullets(self, paragraph, theme_style):
        """Style a single bullet/body paragraph."""
        font = paragraph.font
        font.name  = theme_style["body_font"]
        font.size  = Pt(18)
        font.color.rgb = theme_style["text"]

    # ---------------------------------
    # HELPERS  (internal — do not alter public signatures)
    # ---------------------------------
    def _rgb_to_hex(self, rgb_color):
        return f"{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"

    def _add_solid_rect(self, slide, left, top, width, height, rgb):
        """Add a filled rectangle shape (no outline)."""
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
        shape.line.fill.background()       # no border
        return shape

    def _add_text_box(self, slide, left, top, width, height,
                      text, font_name, font_size, bold, rgb, align=PP_ALIGN.LEFT, italic=False):
        """Add a text box with a single run."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf    = txBox.text_frame
        tf.word_wrap = True
        p    = tf.paragraphs[0]
        p.alignment = align
        run  = p.add_run()
        run.text = text
        run.font.name  = font_name
        run.font.size  = font_size
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = rgb
        return txBox

    # ---------------------------------
    # TITLE SLIDE DECORATOR
    # ---------------------------------
    def _decorate_title_slide(self, slide, topic, theme_style):
        """
        Hero panel title slide.
        Uses theme_style["slide_title"] — a color dedicated solely to the
        presentation topic, separate from content-slide headings.
        """
        PAD = Inches(0.5)

        # Hero panel: left 62 % of slide
        hero_w = Inches(8.2)
        self._add_solid_rect(slide, 0, 0, hero_w, SLIDE_H, theme_style["title_bg"])

        # Thin vertical accent stripe on left edge
        self._add_solid_rect(slide, 0, 0, Inches(0.12), SLIDE_H, theme_style["accent"])

        # Presentation topic — uses slide_title color (its own dedicated color)
        self._add_text_box(
            slide,
            left=Inches(0.55), top=Inches(2.4),
            width=Inches(7.4),  height=Inches(1.8),
            text=topic,
            font_name=theme_style["title_font"],
            font_size=theme_style["title_size"],
            bold=True,
            rgb=theme_style["slide_title"],     # ← dedicated title-slide color
        )

        # Tagline
        self._add_text_box(
            slide,
            left=Inches(0.55), top=Inches(4.4),
            width=Inches(7.0),  height=Inches(0.6),
            text="Generated Presentation",
            font_name=theme_style["body_font"],
            font_size=Pt(15),
            bold=False,
            italic=True,
            rgb=theme_style["subtitle"],
        )

        # Decorative rule under tagline
        rule = slide.shapes.add_shape(
            1,
            Inches(0.55), Inches(5.1),
            Inches(2.8), Inches(0.04)
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = theme_style["accent"]
        rule.line.fill.background()

    # ---------------------------------
    # CONTENT SLIDE DECORATOR
    # ---------------------------------
    def _decorate_content_slide(self, slide, title_text, bullets, theme_style):
        """
        Render a polished two-zone layout:
          • Top header band  — heading color   (theme_style["title"])
          • Bullet cards     — paragraph color (theme_style["text"])
        Works entirely via shapes.
        """
        HEADER_H = Inches(1.3)
        PAD      = Inches(0.4)

        # ── header band (heading zone) ────────────────────────────────
        self._add_solid_rect(slide, 0, 0, SLIDE_W, HEADER_H, theme_style["heading_bg"])

        # Accent stripe across top of header
        self._add_solid_rect(slide, 0, 0, SLIDE_W, Inches(0.07), theme_style["accent"])

        # Heading text — uses "title" color (heading color)
        self._add_text_box(
            slide,
            left=PAD, top=Inches(0.22),
            width=SLIDE_W - PAD * 2, height=HEADER_H - Inches(0.22),
            text=title_text,
            font_name=theme_style["title_font"],
            font_size=theme_style["title_size"],
            bold=True,
            rgb=theme_style["title"],          # ← heading color
        )

        # ── bullet cards (paragraph zone) ────────────────────────────
        card_h        = Inches(0.78)
        card_gap      = Inches(0.18)
        card_start_y  = HEADER_H + Inches(0.35)
        card_left     = PAD
        card_w        = SLIDE_W - PAD * 2
        stripe_w      = Inches(0.07)

        for i, bullet in enumerate(bullets):
            top = card_start_y + i * (card_h + card_gap)

            if top + card_h > SLIDE_H - Inches(0.3):
                break

            # Card background
            self._add_solid_rect(slide, card_left, top, card_w, card_h,
                                  theme_style["card_bg"])

            # Left accent stripe on card
            self._add_solid_rect(slide, card_left, top, stripe_w, card_h,
                                  theme_style["accent"])

            # Paragraph text — uses "text" color (paragraph color)
            self._add_text_box(
                slide,
                left=card_left + stripe_w + Inches(0.2),
                top=top + Inches(0.14),
                width=card_w - stripe_w - Inches(0.3),
                height=card_h - Inches(0.14),
                text=bullet,
                font_name=theme_style["body_font"],
                font_size=theme_style["body_size"],
                bold=False,
                rgb=theme_style["text"],       # ← paragraph color
            )

    # ---------------------------------
    # MAIN FUNCTION
    # ---------------------------------
    def create_ppt(self, topic, content, theme="modern"):

        theme_style = THEMES.get(theme, THEMES["modern"])

        # Generate structured slides
        response = self.think(content)
        data = json.loads(response)

        prs = Presentation()

        # ── set widescreen 16:9 canvas ────────────────────────────────
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        # ══ TITLE SLIDE ═══════════════════════════════════════════════
        title_layout = prs.slide_layouts[6]   # blank layout
        slide = prs.slides.add_slide(title_layout)

        self.apply_theme(slide, theme_style)

        try:
            slide.shapes.title.text = topic
            self.style_title(slide.shapes.title.text_frame, theme_style)
            subtitle = slide.placeholders[1]
            subtitle.text = "Generated Presentation"
            self.style_bullets(subtitle.text_frame.paragraphs[0], theme_style)
        except Exception:
            pass

        self._decorate_title_slide(slide, topic, theme_style)

        # ══ CONTENT SLIDES ════════════════════════════════════════════
        for slide_data in data["slides"]:

            layout = prs.slide_layouts[6]   # blank
            slide  = prs.slides.add_slide(layout)

            self.apply_theme(slide, theme_style)

            # Legacy placeholder fill (kept for compatibility)
            try:
                slide.shapes.title.text = slide_data["title"]
                self.style_title(slide.shapes.title.text_frame, theme_style)

                body = slide.shapes.placeholders[1]
                tf   = body.text_frame
                tf.clear()

                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        tf.text = bullet
                        self.style_bullets(tf.paragraphs[0], theme_style)
                    else:
                        p = tf.add_paragraph()
                        p.text   = bullet
                        p.level  = 0
                        self.style_bullets(p, theme_style)
            except Exception:
                pass

            # Paint the polished card layout on top
            self._decorate_content_slide(
                slide,
                slide_data["title"],
                slide_data["bullets"],
                theme_style,
            )

        return prs